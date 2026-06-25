#!/usr/bin/env python3
"""Build the editorial 'Human Edge' deck as a real, editable PowerPoint.

Cream paper / ink / vermilion; Impact display + Georgia serif + Arial body +
JetBrains Mono labels; hero sketches in bordered hard-shadow frames; colour-block
section dividers; big pull-quotes; and the two structural graphics (trust-tool
matrix, five-differences table) on their own slides.

Editable > bespoke. Regenerate after editing content below:
  python3 scripts/build-deck-pptx.py
  # then re-export the PDF via PowerPoint (File -> Export -> PDF), or the osascript
Output: content/the-human-edge-deck.pptx
"""
from __future__ import annotations
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent
IMG = ROOT / "content" / "images"
OUT = ROOT / "content" / "the-human-edge-deck.pptx"

PAPER, PAPER2 = RGBColor(0xf4, 0xef, 0xe3), RGBColor(0xec, 0xe5, 0xd4)
INK, INKSOFT = RGBColor(0x18, 0x16, 0x13), RGBColor(0x4a, 0x44, 0x3c)
VERM, OCHRE, WHITE = RGBColor(0xe3, 0x4f, 0x29), RGBColor(0xcf, 0x9a, 0x2f), RGBColor(0xff, 0xff, 0xff)
FAINT = RGBColor(0xcb, 0xc3, 0xb0)
FD, FS, FN, FM = "Impact", "Georgia", "Arial", "JetBrains Mono"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = 13.333, 7.5
BLANK = prs.slide_layouts[6]
ML = 0.95
CW = SW - 2 * ML
_N = [0]


def newslide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = bg
    _N[0] += 1
    return s, _N[0]


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def text(s, l, t, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        runs = para if isinstance(para, list) else [para if isinstance(para, tuple) else (para, {})]
        for rtext, st in runs:
            r = p.add_run(); r.text = rtext
            r.font.name = st.get("font", FN)
            r.font.size = Pt(st.get("size", 16))
            r.font.bold = st.get("bold", False)
            r.font.italic = st.get("italic", False)
            r.font.color.rgb = st.get("color", INK)
            if "spacing" in st: p.line_spacing = st["spacing"]
            if st.get("space_after") is not None: p.space_after = Pt(st["space_after"])
    return tb


def bignum(s, n):
    text(s, SW - 3.4, 0.35, 3.0, 2.2, [(f"{n:02d}", dict(font=FD, size=110, color=FAINT))], align=PP_ALIGN.RIGHT)


def kicker(s, label, l=ML, t=0.7, color=VERM):
    text(s, l, t, CW, 0.4, [(label.upper(), dict(font=FM, size=11, color=color, bold=True))])


def img_fit(name, maxw, maxh):
    iw, ih = Image.open(IMG / name).size
    sc = min(maxw / iw, maxh / ih)
    return iw * sc, ih * sc


def img_frame(s, name, l, t, w, h):
    rect(s, l + 0.14, t + 0.14, w, h, fill=INK)              # hard shadow
    rect(s, l, t, w, h, fill=PAPER2, line=INK, lw=1.75)      # matted frame
    iw, ih = img_fit(name, w - 0.28, h - 0.28)               # fit, preserve aspect
    pic = s.shapes.add_picture(str(IMG / name), Inches(l + (w - iw) / 2), Inches(t + (h - ih) / 2), Inches(iw), Inches(ih))
    pic.shadow.inherit = False
    return pic


# ---------- slide builders ----------
def cover():
    s, n = newslide(); bignum(s, n)
    rect(s, ML, 1.0, 2.7, 0.45, fill=VERM)
    text(s, ML, 1.0, 4.2, 0.45, [("ONE-DAY MASTERCLASS", dict(font=FM, size=11, color=WHITE, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    kicker(s, "Curtin Executive Education · Dr Michael Borck", t=1.75)
    text(s, ML, 2.5, 11.4, 4.2,
         [("THE", dict(font=FD, size=84, color=INK, spacing=0.92)),
          ("HUMAN", dict(font=FD, size=84, color=INK, spacing=0.92)),
          ("EDGE", dict(font=FD, size=84, color=INK, spacing=0.92))])
    text(s, ML, 6.7, 11.4, 0.5, [("Using & delivering AI with judgement.", dict(font=FS, size=24, italic=True, color=INKSOFT))])


def section(bg, eyebrow, title_lines, lead):
    s, n = newslide(bg)
    text(s, ML, 1.7, 11.4, 0.5, [(eyebrow.upper(), dict(font=FM, size=14, color=WHITE, bold=True))])
    text(s, ML, 2.5, 11.8, 3.8, [(ln, dict(font=FD, size=64, color=WHITE, spacing=0.95)) for ln in title_lines])
    text(s, ML, SH - 1.7, 9.8, 1.2, [(lead, dict(font=FS, size=22, italic=True, color=WHITE))])


def quote(label, qtext, attr):
    s, n = newslide(); bignum(s, n)
    kicker(s, label)
    text(s, ML, 1.9, 1.5, 1.6, [("“", dict(font=FS, size=180, color=VERM))])
    text(s, ML + 0.6, 2.7, 11.4, 3.2, [(qtext, dict(font=FS, size=34, color=INK, spacing=1.06))])
    text(s, ML + 0.6, SH - 1.3, 11.4, 0.5, [(attr.upper(), dict(font=FM, size=11, color=INKSOFT, bold=True))])


def concept(label, title, body, img_name):
    s, n = newslide(); bignum(s, n)
    img_frame(s, img_name, ML, 1.4, 6.0, 5.0)
    tx, tw = 7.5, 5.0
    text(s, tx, 1.5, tw, 0.4, [(label.upper(), dict(font=FM, size=11, color=VERM, bold=True))])
    text(s, tx, 2.0, tw, 2.4, [(title, dict(font=FS, size=34, color=INK, bold=True, spacing=1.0))])
    text(s, tx, 4.4, tw, 2.6, [(body, dict(font=FN, size=14.5, color=INKSOFT, spacing=1.22))])


def grid(label, title, title_tail, items, cols=2, footleft=None):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    tline = [(title, dict(font=FS, size=40, color=INK, bold=True))]
    if title_tail: tline.append((" " + title_tail, dict(font=FS, size=40, color=INK, italic=True)))
    text(s, ML, 1.5, CW, 1.6, [tline])
    top = 3.5; colw = (CW - (cols - 1) * 0.6) / cols
    for i, (lab, desc) in enumerate(items):
        c, r = i % cols, i // cols
        x = ML + c * (colw + 0.6); y = top + r * 1.85
        rect(s, x, y, colw, 0.03, fill=INK)
        text(s, x, y + 0.12, colw, 0.6, [(lab, dict(font=FD, size=20, color=VERM))])
        text(s, x, y + 0.72, colw, 1.0, [(desc, dict(font=FN, size=12.5, color=INKSOFT, spacing=1.12))])
    if footleft: text(s, ML, SH - 0.55, CW, 0.3, [(footleft, dict(font=FM, size=9, color=INKSOFT))])


def feature(label, title, body, footleft=None):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    text(s, ML, 1.7, 10.8, 2.2, [(title, dict(font=FS, size=40, color=INK, bold=True, spacing=1.0))])
    text(s, ML, 4.0, 10.8, 3.0, [(body, dict(font=FN, size=15.5, color=INKSOFT, spacing=1.28))])
    if footleft: text(s, ML, SH - 0.55, CW, 0.3, [(footleft, dict(font=FM, size=9, color=INKSOFT))])


def listslide(label, title, items, intro=None):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    text(s, ML, 1.5, CW, 1.3, [(title, dict(font=FS, size=40, color=INK, bold=True, spacing=1.0))])
    top = 3.3
    paras = []
    if intro:
        paras.append([(intro, dict(font=FN, size=15, color=INKSOFT, italic=True))])
    for it in items:
        paras.append([("▸  ", dict(font=FN, size=16, color=VERM, bold=True)), (it, dict(font=FN, size=15.5, color=INK, spacing=1.15))])
    text(s, ML, top, CW, 4.0, paras)


def steps(label, title, items):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    text(s, ML, 1.5, CW, 1.2, [(title, dict(font=FS, size=40, color=INK, bold=True))])
    top = 3.3; colw = (CW - 2 * 0.7) / 3
    for i, (num, h, body) in enumerate(items):
        x = ML + i * (colw + 0.7)
        rect(s, x, top, colw, 0.04, fill=INK)
        text(s, x, top + 0.15, colw, 0.9, [(num, dict(font=FD, size=32, color=VERM))])
        text(s, x, top + 1.0, colw, 0.8, [(h, dict(font=FS, size=18, color=INK, bold=True))])
        text(s, x, top + 1.8, colw, 2.0, [(body, dict(font=FN, size=12.5, color=INKSOFT, spacing=1.15))])


def sprint(num, title, who, out, body):
    s, n = newslide(); bignum(s, n)
    text(s, ML, 1.4, 2.4, 1.6, [(num, dict(font=FD, size=120, color=VERM, spacing=0.9))])
    text(s, ML + 2.6, 1.7, CW - 2.6, 1.4, [(title, dict(font=FS, size=40, color=INK, bold=True, spacing=0.98))])
    text(s, ML + 2.6, 3.4, CW - 2.6, 1.8, [(body, dict(font=FN, size=16, color=INKSOFT, spacing=1.3))])
    text(s, ML + 2.6, SH - 1.7, CW - 2.6, 0.5, [("INTERVIEW  " + who, dict(font=FM, size=11, color=INKSOFT, bold=True))])
    text(s, ML + 2.6, SH - 1.15, CW - 2.6, 0.6, [("→ " + out, dict(font=FS, size=17, italic=True, color=VERM, bold=True))])



def brk(img_name, time, label, back):
    s, n = newslide(); bignum(s, n)
    rect(s, SW/2 - 1.7, 1.7 + 0.14, 3.4, 2.5, fill=INK)
    pic = s.shapes.add_picture(str(IMG / img_name), Inches(SW/2 - 1.7), Inches(1.7), Inches(3.4), Inches(2.5))
    pic.line.color.rgb = INK; pic.line.width = Pt(1.75); pic.shadow.inherit = False
    text(s, ML, 4.5, CW, 1.3, [(time, dict(font=FD, size=84, color=VERM))], align=PP_ALIGN.CENTER)
    text(s, ML, 5.8, CW, 0.8, [(label, dict(font=FS, size=32, italic=True, color=INK))], align=PP_ALIGN.CENTER)
    text(s, ML, SH - 0.85, CW, 0.4, [(back.upper(), dict(font=FM, size=11, color=INKSOFT, bold=True))], align=PP_ALIGN.CENTER)


def resources(label, title, items):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    text(s, ML, 1.5, CW, 1.0, [(title, dict(font=FS, size=38, color=INK, bold=True))])
    cols = len(items); gap = 2.2 if cols <= 2 else 1.4
    colw = (CW - (cols - 1) * gap) / cols; qsize = min(2.5, colw - 0.25); qtop = 2.6
    for i, (ttl, url, qr, note) in enumerate(items):
        x = ML + i * (colw + gap); qx = x + (colw - qsize) / 2
        rect(s, qx + 0.12, qtop + 0.12, qsize, qsize, fill=INK)
        rect(s, qx, qtop, qsize, qsize, fill=WHITE, line=INK, lw=1.5)
        iw, ih = img_fit(qr, qsize - 0.4, qsize - 0.4)
        pic = s.shapes.add_picture(str(IMG / qr), Inches(qx + (qsize - iw) / 2), Inches(qtop + (qsize - ih) / 2), Inches(iw), Inches(ih))
        pic.shadow.inherit = False
        disp = url.replace("https://", "").replace("http://", "").rstrip("/")
        text(s, x, qtop + qsize + 0.18, colw, 0.5, [(ttl, dict(font=FS, size=19, color=INK, bold=True))], align=PP_ALIGN.CENTER)
        text(s, x, qtop + qsize + 0.62, colw, 0.45, [(disp, dict(font=FM, size=11, color=VERM, bold=True))], align=PP_ALIGN.CENTER)
        text(s, x, qtop + qsize + 1.05, colw, 0.75, [(note, dict(font=FS, size=11.5, italic=True, color=INKSOFT))], align=PP_ALIGN.CENTER)

def schedule(title, rows):
    s, n = newslide(); bignum(s, n); kicker(s, "The day at a glance")
    text(s, ML, 1.5, CW, 1.0, [(title, dict(font=FS, size=38, color=INK, bold=True))])
    top = 2.9; rh = 0.52
    for i, (t, b) in enumerate(rows):
        y = top + i * rh
        text(s, ML, y, 2.7, rh, [(t, dict(font=FM, size=12.5, color=VERM, bold=True))])
        text(s, ML + 2.9, y, CW - 2.9, rh, [(b, dict(font=FN, size=13.5, color=INK))])

def activity(minutes, title, brief):
    s, n = newslide()
    rect(s, 0, 0, SW, 0.62, fill=VERM)
    text(s, ML, 0, 5.0, 0.62, [("ACTIVITY", dict(font=FM, size=13, color=WHITE, bold=True))], anchor=MSO_ANCHOR.MIDDLE)
    text(s, SW - 4.5, 0, 3.7, 0.62, [("≈ " + minutes, dict(font=FM, size=13, color=WHITE, bold=True))], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML, 2.5, CW, 1.8, [(title, dict(font=FS, size=38, color=INK, bold=True, spacing=0.98))])
    text(s, ML, 4.4, CW, 1.8, [(brief, dict(font=FN, size=16, color=INKSOFT, spacing=1.3))])

def wide(label, title, img_name, caption):
    s, n = newslide(); bignum(s, n); kicker(s, label)
    text(s, ML, 1.4, CW, 1.0, [(title, dict(font=FS, size=38, color=INK, bold=True))])
    iw, ih = img_fit(img_name, 11.6, 3.5)
    l = (SW - iw) / 2; t = 3.7
    rect(s, l + 0.12, t + 0.12, iw, ih, fill=INK)
    pic = s.shapes.add_picture(str(IMG / img_name), Inches(l), Inches(t), Inches(iw), Inches(ih))
    pic.line.color.rgb = INK; pic.line.width = Pt(1.5); pic.shadow.inherit = False
    text(s, ML, t + ih + 0.32, CW, 0.5, [(caption, dict(font=FM, size=11, color=INKSOFT))], align=PP_ALIGN.CENTER)

def close_deck():
    s, n = newslide(); bignum(s, n)
    text(s, ML, 2.1, CW, 0.5, [("THE AI IS YOUR STARTING POINT, NOT YOUR FINISH LINE", dict(font=FM, size=12, color=VERM, bold=True))], align=PP_ALIGN.CENTER)
    text(s, ML, 2.9, CW, 3.2,
         [("STAY THE ONE", dict(font=FD, size=52, color=INK, spacing=0.95)),
          ("HOLDING THE", dict(font=FD, size=52, color=INK, spacing=0.95)),
          ("JUDGEMENT.", dict(font=FD, size=52, color=INK, spacing=0.95))], align=PP_ALIGN.CENTER)
    text(s, ML, SH - 0.9, CW, 0.5, [("michael-borck.github.io/the-human-edge", dict(font=FS, size=18, italic=True, color=INKSOFT))], align=PP_ALIGN.CENTER)


# ---------- build ----------
cover()
resources("Before we begin — get connected", "Check you can open all three.", [
    ("Free chat", "https://chat.locopuente.org", "qr-free-chat.png", "Today only — gone after the course."),
    ("Companion site", "https://michael-borck.github.io/the-human-edge/", "qr-companion-site.png", "Yours to keep — before, during, after."),
    ("RetailFlow portal", "https://retailflow.eduserver.au", "qr-retailflow.png", "This afternoon's sprints · access code: pilot2024"),
])
schedule("The day.", [
    ("9:00–10:30", "The average, the tool, the edge — foundations + the trust tool"),
    ("10:30–11:00", "Morning tea"),
    ("11:00–12:30", "Using AI well — RTCF, the two-pass demo, workflow redesign"),
    ("12:30–1:15", "Lunch"),
    ("1:15–2:30", "Why AI delivery is different — the five differences"),
    ("2:30–3:00", "Afternoon tea"),
    ("3:00–4:00", "Designing the human in — the three sprints + go/no-go"),
    ("4:00–4:30", "The edge that's left to humans + your action plan"),
])
quote("The provocation we'll spend the day on",
      "If the AI is good at running your work, it's good at running everyone else's. Generic competence is the baseline — not the advantage.",
      "So where does your edge live?")
grid("Today", "One idea. Two scales.", None,
     [("MORNING", "The judgement, at the level of your own work: the trust tool, RTCF, the two-pass demo."),
      ("AFTERNOON", "The same judgement, lifted to delivery: the five differences, human-in-the-loop, the go/no-go.")],
     footleft="The trust tool — taught once, reused at project scale")
section(VERM, "Part one", ["THE AVERAGE,", "THE TOOL,", "YOUR EDGE."],
        "Personal fluency — the judgement that lasts after the next model ships.")
grid("A quick note on scope", "AI is bigger than this.", "Today is one slice.",
     [("THE WHOLE FIELD", "Machine learning, computer vision, audio, CNNs, robotics, forecasting — decades of work, far beyond chatbots."),
      ("TODAY'S FOCUS", "Generative AI — models that produce text, images and code. The part reshaping everyday work right now, and the part you'll use at your desk.")],
     footleft="The judgement we build here transfers across the whole field")
concept("What AI actually is", "It predicts the next word.",
        "Autocomplete, scaled up. Trained on billions of pages, it learned what tends to follow what. The result looks like understanding. It's pattern-matching.",
        "autocomplete-on-steroids.jpg")
concept("The one mental model to keep", "A convincing average.",
        "Fluent, plausible, usually about-right. But \"about right\" and \"exactly right\" are not the same — and knowing which you need is the whole game.",
        "convincing-average.jpg")
listslide("Why it's brilliant at some things", "Brilliant here.",
          ["Drafting and writing", "Summarising documents", "Formatting and transforming",
           "Brainstorming", "Explaining concepts in simpler terms", "Finding patterns in data"],
          intro="Tasks where \"the most plausible version\" is exactly what you want.")
listslide("…and quietly unreliable at others", "Quietly wrong there.",
          ["Factual accuracy — plausible text, not verified truth", "Precise figures and calculations",
           "Multi-step logical reasoning", "Your context — it doesn't know your company or politics",
           "Ethical judgement"],
          intro="Tasks that need a single correct answer it can't verify.")
concept("Not a bug — a property", "Confident. Wrong. Same polish.",
        "It never says \"I don't know.\" The fluent-but-wrong answer looks identical to the fluent-and-right one. Nothing in the tone warns you.",
        "confident-nonsense.jpg")
concept("The spine of the day", "The trust tool.",
        "Two questions, four corners — lean in where 'about right' is fine; keep a human in the loop where exactly-right meets high stakes.",
        "trust-tool-matrix.png")
grid("Two questions, four corners", "Average or precise? Small or large?", None,
     [("AVERAGE + SMALL", "Lean in. AI's home turf — let it run."),
      ("AVERAGE + LARGE", "Lean in, with a glance. Sanity-check before it goes out."),
      ("PRECISE + SMALL", "Use, then verify. Confirm the specifics."),
      ("PRECISE + LARGE", "Human in the loop. A person owns the decision.")],
     cols=2)
feature("The trust tool, on real tasks", "Average or precise? Small or large?",
        "Drafting an internal email. Average is fine, stakes small — lean in.\nBoard-paper figures. Precision required, stakes large — a human owns every number.\nSummarising fifty complaints for themes. Lean in for the themes, then skim for the one complaint the AI smoothed away.")
activity("20 min", "Exercise 1 — the AI tool test drive",
         "Pick a real, small task. Run it through an AI tool. Judge the output with the trust tool — average or precise? small or large stakes? — and spot the confidently-wrong bit. (Worksheet: in your workbook.)")

feature("Exercise 1 — the AI tool test drive", "Same task. Judge the result with the trust tool.",
        "Pick a real, small task from your work — an email, a summary, a first draft. Run it through an AI tool. Look at the output: where on the grid does this task sit? Is it 'about right' or does it need exactness? Where would it be confidently wrong? Share one place it was useful, one place it was off. There's no 'best' tool — never assume the first output is ready to use.")
brk("break-morning-tea.jpg", "10:30", "Morning tea", "Back at 11:00")
activity("25 min", "Build your prompt library (RTCF)",
         "For three real tasks: write the naive version, then rebuild with RTCF. Context is where your edge enters. (Worksheet: Your prompt library.)")

grid("Talking to AI well", "RTCF — two halves.", None,
     [("R · T · F  SCAFFOLDING", "Weaker & self-hosted models need it spelled out. A frontier model infers it from clear writing."),
      ("C  IS YOUR JUDGEMENT", "We call it Context — but it's really your taste, experience & critical thinking. The one thing no model supplies. So make yours different.")],
     footleft="The durable skill: being clear about what you want")
feature("Before & after", "Unusable. Then sendable.",
        "Without: \"Write an email about the project delay.\"\nWith RTCF: a senior PM, professional & empathetic, writes a client email explaining a 2-week delay (vendor integration; client values transparency; new date March 15), under 200 words — acknowledge → explain → new date → recommit.\nWe call it Context — really, it's your judgement going in.")
activity("15 min", "The two-pass demo — live",
         "One task, run twice. Pass 1: naive → slick, identical for everyone. Pass 2: add your edge → it gets a soul. Let the room shout the edge inputs.")

concept("The money moment — live", "Add your edge. It gets a soul.",
        "Pass 1: a naive prompt → slick, correct, identical for everyone. Pass 2: layer in what only you hold. Watch the output turn yours.",
        "two-pass-demo.jpg")
activity("25 min", "Exercise 2 — redesign your workflow",
         "Pick a task you own. Map its steps, judge each with the trust tool, redesign so AI handles the average and you keep the high-stakes parts. (Worksheet: Redesign your workflow.)")

feature("Exercise 2 — redesign your workflow", "Redesign one workflow.",
        "Pick a task you own. List its steps. Judge each with the trust tool — where does AI genuinely help, where must your judgement stay in charge? Redesign so AI handles the average and you keep the high-stakes parts. You leave with something to trial next week.")
feature("Where we are by lunch", "A mental model, a grid, a method.",
        "AI produces a convincing average. The trust tool tells you when to lean in. RTCF directs it well — and Context is the part that survives any model. One question for the afternoon: if this is how I trust AI on a single task, how do I trust it across a whole project?")
brk("break-lunch.jpg", "12:30", "Lunch", "Back at 1:15")
section(INK, "Part two", ["WHY AI", "DELIVERY IS", "DIFFERENT."],
        "You're the delivery lead. Ship one funded initiative — without the predictable failures.")
feature("A quick reminder before we go on", "This applies beyond text.",
        "This morning was generative AI. But the five differences aren't LLM-specific — they describe any AI project. A vision system is non-deterministic too; an audio model's data is discovered, not ordered; any AI output — image, sound, prediction — needs a human checkpoint where the stakes are high. The judgement you build this afternoon travels across the whole field: CNNs, audio, ML, all of it.")
quote("The fair question",
      "I've delivered projects before. What's actually different about an AI project?",
      "They don't fail more often. They fail differently — in places you didn't see coming.")
concept("Five ways AI breaks the rules you usually lead by", "They fail differently.",
        "For each assumption a normal project relies on — why AI breaks it, and the leadership move it forces.",
        "five-differences-summary.png")
concept("Difference 01", "\"Done\" won't hold still.",
        "Ask twice, get two answers. There's no single correct output — only a spread. Stop chasing a fixed spec; decide what \"good enough\" looks like for this use.",
        "done-cant-be-specified.jpg")
concept("Difference 02", "The demo is a trap.",
        "It shows the happy path. The messy real-world cases are where it falls over. Treat the demo as the start of the hard part — not the end.",
        "demo-is-a-trap.jpg")
concept("Difference 03", "The data is the uncertainty.",
        "You don't order it; you discover it — its gaps, biases and surprises show up only when you work with it. Run delivery as discovery, and build the learning into the plan.",
        "data-is-the-uncertainty.jpg")
concept("Difference 04", "Verification is the product.",
        "Outputs vary and can be confidently wrong. The value lives in how you check them, and where a human stays in the loop. Design the checking deliberately — make it part of what you ship.",
        "verification-is-the-product.jpg")
concept("Difference 05", "Generic is the baseline, not the edge.",
        "When every rival runs the same models, the model isn't your advantage — anyone can have it. Compete on the human variation a rival's identical model can't reproduce.",
        "where-your-edge-lives.jpg")
quote("The thread through all five",
      "AI gives you something fluent and roughly-right, fast — then asks you to decide when roughly-right is good enough, and who checks when it isn't.",
      "A leadership question. Which is why it lands on your desk.")
feature("Your project · RetailFlow", "You're the delivery lead.",
        "The board funded one AI initiative. Ship it without the predictable failures. Carry it through three moves — scoping and stress-testing it by interviewing the RetailFlow team: live chatbots with opinions, who disagree, and won't do your job for you.   access code: pilot2024")
activity("45 min", "Sprint 1 — scope it against reality",
         "You're the delivery lead for one RetailFlow initiative. Interview Priya & Marcus. Produce a scoped objective, data needs, the trust-tool split, and a defined 'good enough'. (Worksheet: Sprint 1. Access code: pilot2024.)")

sprint("01", "Scope it against reality", "Priya (Data) & Marcus (CIO)",
       "Scoped objectives, data requirements, a defined 'good enough'",
       "Reconcile 'move fast' with 'the data isn't ready.' Marcus pushes speed; Priya gives the honest, data-grounded timeline. Your scope is where you reconcile the two — and you'll defend it later.")
brk("break-afternoon-tea.jpg", "2:30", "Afternoon tea", "Back at 3:00")
activity("45 min", "Sprint 2 — stakeholders & human-in-the-loop",
         "Continue your initiative. Interview Emma, Tom & David. Build a stakeholder plan and design the human-in-the-loop checkpoints — each named, marked builds-vs-consumes.")

sprint("02", "Stakeholders & human-in-the-loop", "Emma (MD), Tom (frontline) & David (CFO)",
       "A stakeholder plan and a human-in-the-loop checkpoint design",
       "Decide where a human must stay in charge. This is the trust tool and difference 4 in action: design the checking deliberately — where does verification sit, where must a person own the decision?")
activity("45 min", "Sprint 3 — roadmap, risk & the go/no-go",
         "Build the delivery roadmap with gates and a risk register. Then make the call — Scale, Pivot, or Kill — and defend it. (Crisis element if time allows.)")

sprint("03", "Roadmap, risk & the go/no-go", "Build it, then defend it",
       "A delivery roadmap with gates, a risk register, the Scale / Pivot / Kill call",
       "Build the roadmap with go/no-go gates and a risk register. Then make the call: Scale, Pivot, or Kill. The question isn't 'is it perfect?' — it's whether the evidence justifies continuing.")
concept("The go / no-go call", "Scale. Pivot. Kill.",
        "At each gate the question isn't \"is it perfect?\" — it's whether the evidence justifies continuing. Killing early is success, not failure.",
        "scale-pivot-kill.jpg")
grid("What you've built by 4:00", "A delivery design, not a theory.", None,
     [("SCOPE", "Objectives, data requirements, a defined \"good enough.\""),
      ("ROADMAP", "Milestones and go/no-go gates."),
      ("CHECKPOINTS", "A stakeholder plan and human-in-the-loop design."),
      ("RISK", "A register and a Scale / Pivot / Kill call you can defend.")])
concept("Back to the provocation", "The edge is your judgement.",
        "It looks like all we're doing is adding context. We're not — we're adding your judgement: taste, experience, critical thinking. Averages converge; yours diverges. Make yours different — the generic version is the one your competitors already have.",
        "where-your-edge-lives.jpg")
grid("What you do next", "Your action plan.", None,
     [("WEEK", "Trial your redesigned workflow. Run one task through the trust tool first."),
      ("MONTH", "Make RTCF a habit. Start one conversation about where a human stays in the loop."),
      ("QUARTER", "Scope one real initiative — and defend the Scale / Pivot / Kill call.")],
     cols=3)
grid("What you leave with", "Tools that don't expire.", None,
     [("DECISION TOOLS", "The trust tool, RTCF, the five differences, Scale / Pivot / Kill — any platform, any model."),
      ("YOUR WORK", "A redesigned workflow and a scoped project design you can act on."),
      ("A PLAN", "Specific commitments for the week, month and quarter."),
      ("GO DEEPER", "michael-borck.github.io/the-human-edge — incl. the two-page voice method.")])
wide("Go deeper — the shape underneath it all", "The conversation loop.", "conversation-loop.png",
     "Brainstorm → Ideate → Iterate → Amplify → Repeat. You stay in the conversation.")

grid("Go deeper — push back on AI output", "VET: verify, explain, test.", None,
     [("VERIFY", "Can I find this independently? Check sources, cross-reference, look up citations — AI fabricates references routinely."),
      ("EXPLAIN", "Can I say this in my own words? If you can't rephrase it, you're holding words, not understanding."),
      ("TEST", "Does it hold up? Play devil's advocate, check the edge cases, ask 'what if?'"),],
     cols=3, footleft="The trust tool says when to check; VET says how.  companion site → Go deeper")
feature("Go deeper — the two-chat workflow", "Separate thinking from building.",
        "Session 1 — Thinking Chat: explore, challenge, follow tangents, find your real question (messy by design). Then stop: read what emerged, decide what matters, write a focused brief. Session 2 — Build Chat: produce the output from that brief. You don't copy-paste between them — that curation is where your judgement lives. Neither session has the full picture; only you hold both.")

close_deck()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f">> wrote {OUT.relative_to(ROOT)}  ({_N[0]} slides)")
